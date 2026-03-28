"""
HR Policy Knowledge Vectorizer.

Ingests policy documents from data/knowledge/ (or sample_knowledge/),
generates embeddings via Azure AI Foundry, and persists them in a local
ChromaDB collection. Run once at setup; re-run to refresh after policy updates.

Supported formats: .txt, .md, .pdf, .doc, .docx, .pptx

Usage:
    uv run hr-seed                                   # seed from sample_knowledge/
    uv run hr-seed --source data/knowledge/          # seed from data/knowledge/
    uv run hr-seed --reset                           # clear and re-seed
"""

from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

from rich.console import Console
from rich.progress import Progress, SpinnerColumn, TextColumn

console = Console()

COLLECTION_NAME = "hr_policy"
CHUNK_SIZE = 800       # characters per chunk
CHUNK_OVERLAP = 100    # character overlap between chunks

# All supported file extensions
SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".doc", ".docx", ".pptx"}

# Files that should never be indexed — not HR policy content
EXCLUDED_FILENAMES = {
    "Contoso-HR-Policy.doc",          # old binary .doc — replaced by -v2.docx
    "Copilot-Studio-HR-Scenario.pptx", # product deployment guide, not HR policy
}


def extract_text(file_path: Path) -> str:
    """Extract plain text from any supported document format.

    Args:
        file_path: Path to the source document.

    Returns:
        Extracted plain text, or empty string if extraction fails.
    """
    suffix = file_path.suffix.lower()

    if suffix in (".txt", ".md"):
        try:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        except Exception as e:
            console.print(f"[yellow]  Warning: could not read {file_path.name}: {e}[/]")
            return ""

    if suffix == ".pdf":
        return _extract_pdf(file_path)

    if suffix in (".doc", ".docx"):
        return _extract_doc(file_path)

    if suffix == ".pptx":
        return _extract_pptx(file_path)

    return ""


def _extract_pdf(file_path: Path) -> str:
    """Extract text from a PDF file using pypdf."""
    try:
        import pypdf
        reader = pypdf.PdfReader(str(file_path))
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                parts.append(text.strip())
        return "\n\n".join(parts)
    except ImportError:
        console.print(
            "[yellow]  pypdf not installed — skipping PDF. Run: uv add pypdf[/]"
        )
        return ""
    except Exception as e:
        console.print(f"[yellow]  Warning: could not extract PDF {file_path.name}: {e}[/]")
        return ""


def _extract_doc(file_path: Path) -> str:
    """Extract text from .doc / .docx files using python-docx."""
    try:
        import docx
        doc = docx.Document(str(file_path))
        parts = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(parts)
    except ImportError:
        console.print(
            "[yellow]  python-docx not installed — skipping DOC/DOCX. Run: uv add python-docx[/]"
        )
        return ""
    except Exception as e:
        console.print(
            f"[yellow]  Warning: could not extract DOC {file_path.name}: {e}[/]"
        )
        # .doc (old binary format) requires extra handling; fall back gracefully
        if file_path.suffix.lower() == ".doc":
            console.print(
                f"[dim]  Note: .doc (binary Word) requires LibreOffice or antiword "
                f"for extraction. Convert to .docx for best results.[/]"
            )
        return ""


def _extract_pptx(file_path: Path) -> str:
    """Extract text from .pptx files using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except ImportError:
        console.print(
            "[yellow]  python-pptx not installed — skipping PPTX. Run: uv add python-pptx[/]"
        )
        return ""
    except Exception as e:
        console.print(
            f"[yellow]  Warning: could not extract PPTX {file_path.name}: {e}[/]"
        )
        return ""


def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks for better retrieval.

    Args:
        text: Source text to chunk.
        size: Target chunk size in characters.
        overlap: Overlap between consecutive chunks.

    Returns:
        List of non-empty text chunks.
    """
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = end - overlap
    return chunks


def seed_knowledge(
    source_dir: Path,
    chroma_dir: Path,
    dest_dir: Path,
    embeddings,
    reset: bool = False,
) -> int:
    """Ingest policy documents into ChromaDB.

    Supports .txt, .md, .pdf, .doc, .docx, .pptx files.

    Args:
        source_dir: Directory containing source documents.
        chroma_dir: ChromaDB persistence directory.
        dest_dir: data/knowledge/ — copy files here after seeding.
        embeddings: LangChain-compatible embeddings instance.
        reset: If True, clear existing ChromaDB collection first.

    Returns:
        Number of chunks indexed.
    """
    import chromadb
    from chromadb.utils.embedding_functions import EmbeddingFunction

    class LangChainEmbeddingWrapper(EmbeddingFunction):
        def __init__(self, lc_embeddings):
            self._emb = lc_embeddings

        def __call__(self, input: list[str]) -> list[list[float]]:  # noqa: A002
            return self._emb.embed_documents(input)

    chroma_dir.mkdir(parents=True, exist_ok=True)
    dest_dir.mkdir(parents=True, exist_ok=True)

    client = chromadb.PersistentClient(path=str(chroma_dir))

    if reset:
        try:
            client.delete_collection(COLLECTION_NAME)
            console.print(f"[yellow]Cleared existing collection '{COLLECTION_NAME}'[/]")
        except Exception:
            pass

    ef = LangChainEmbeddingWrapper(embeddings)
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=ef,
        metadata={"description": "Contoso HR policy documents"},
    )

    # Gather all supported files
    files: list[Path] = []
    for ext in sorted(SUPPORTED_EXTENSIONS):
        files.extend(sorted(source_dir.glob(f"*{ext}")))
    # Deduplicate (glob order can vary)
    seen: set[str] = set()
    unique_files: list[Path] = []
    for f in files:
        if f.name not in seen:
            seen.add(f.name)
            unique_files.append(f)

    if not unique_files:
        console.print(f"[yellow]No supported documents found in {source_dir}[/]")
        console.print(f"[dim]Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}[/]")
        return 0

    # Filter out explicitly excluded files
    excluded = [f for f in unique_files if f.name in EXCLUDED_FILENAMES]
    unique_files = [f for f in unique_files if f.name not in EXCLUDED_FILENAMES]
    for f in excluded:
        console.print(f"[dim]  Skipped (excluded): {f.name}[/]")

    console.print(f"[dim]Found {len(unique_files)} documents to ingest[/]")

    total_chunks = 0
    skipped = 0

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        task = progress.add_task("Vectorizing documents...", total=len(unique_files))

        for doc_path in unique_files:
            progress.update(task, description=f"Processing {doc_path.name}...")

            text = extract_text(doc_path)
            if not text.strip():
                console.print(
                    f"[dim]  Skipped {doc_path.name} (no extractable text)[/]"
                )
                skipped += 1
                progress.advance(task)
                continue

            chunks = chunk_text(text)
            if not chunks:
                skipped += 1
                progress.advance(task)
                continue

            # Sanitize stem for use as ChromaDB ID prefix
            safe_stem = "".join(c if c.isalnum() or c in "-_" else "_" for c in doc_path.stem)
            ids = [f"{safe_stem}_{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    "source": doc_path.name,
                    "file_type": doc_path.suffix.lower(),
                    "chunk_index": i,
                }
                for i in range(len(chunks))
            ]

            # Upsert in batches of 50 to stay within API rate limits
            batch_size = 50
            for batch_start in range(0, len(chunks), batch_size):
                batch_end = min(batch_start + batch_size, len(chunks))
                collection.upsert(
                    documents=chunks[batch_start:batch_end],
                    ids=ids[batch_start:batch_end],
                    metadatas=metadatas[batch_start:batch_end],
                )

            total_chunks += len(chunks)

            # Copy to data/knowledge/ for reference
            shutil.copy2(doc_path, dest_dir / doc_path.name)
            progress.advance(task)

    if skipped:
        console.print(f"[dim]Skipped {skipped} file(s) with no extractable text[/]")

    return total_chunks


def main() -> None:
    """CLI entry point for seeding the knowledge base."""
    parser = argparse.ArgumentParser(
        description="Seed Contoso HR policy knowledge base into ChromaDB"
    )
    parser.add_argument(
        "--source",
        type=Path,
        default=None,
        help="Source directory of policy docs (default: sample_knowledge/)",
    )
    parser.add_argument(
        "--reset",
        action="store_true",
        help="Clear existing ChromaDB collection before seeding",
    )
    args = parser.parse_args()

    # Find project root
    here = Path(__file__).resolve()
    for parent in here.parents:
        if (parent / "pyproject.toml").exists():
            project_root = parent
            break
    else:
        project_root = Path.cwd()

    source_dir = args.source or (project_root / "sample_knowledge")
    if not source_dir.exists():
        console.print(f"[red]Source directory not found: {source_dir}[/]")
        sys.exit(1)

    from contoso_hr.config import Config
    config = Config.from_env(project_root)

    errors = config.validate()
    if errors:
        console.print("[red]Configuration errors:[/]")
        for e in errors:
            console.print(f"  [red]• {e}[/]")
        sys.exit(1)

    embeddings = config.get_embeddings()

    console.print(f"[bold cyan]Seeding knowledge base from {source_dir}[/]")
    count = seed_knowledge(
        source_dir=source_dir,
        chroma_dir=config.chroma_dir,
        dest_dir=config.knowledge_dir,
        embeddings=embeddings,
        reset=args.reset,
    )
    console.print(f"[bold green]✓ Indexed {count} chunks into ChromaDB[/]")
